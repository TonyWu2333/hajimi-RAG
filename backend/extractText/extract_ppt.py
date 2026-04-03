"""
extract_ppt.py
提取 PowerPoint 文件（.pptx / .ppt）中的纯文本。

依赖:
    pip install python-pptx
    .ppt 格式还需要安装 LibreOffice（用于格式转换）

用法:
    from extract_ppt import extract_ppt
    text = extract_ppt("path/to/file.pptx")
"""

import subprocess
import tempfile
import shutil
from pathlib import Path


def _convert_ppt_to_pptx(ppt_path: Path) -> Path:
    """使用 LibreOffice 将旧版 .ppt 转换为 .pptx，返回临时文件路径。"""
    with tempfile.TemporaryDirectory() as tmp_dir:
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "pptx",
                "--outdir", tmp_dir,
                str(ppt_path),
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice 转换失败: {result.stderr.strip()}"
            )
        converted = Path(tmp_dir) / (ppt_path.stem + ".pptx")
        if not converted.exists():
            raise RuntimeError("转换后文件未找到，请确认 LibreOffice 已安装。")
        dest = Path(tempfile.mktemp(suffix=".pptx"))
        shutil.copy2(converted, dest)
        return dest


def extract_ppt(file_path: str) -> str:
    """
    提取 PowerPoint 文件的文本内容。

    按幻灯片顺序提取所有形状中的文本框内容，
    每张幻灯片前加 "--- Slide N ---" 标题行。

    Args:
        file_path: PPT/PPTX 文件路径。

    Returns:
        提取的纯文本字符串。

    Raises:
        FileNotFoundError: 文件不存在。
        ImportError: 未安装 python-pptx。
        ValueError: 不支持的文件格式。
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请先安装依赖: pip install python-pptx")

    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    suffix = path.suffix.lower()
    if suffix not in (".pptx", ".ppt"):
        raise ValueError(f"不支持的文件格式: {suffix}，仅支持 .pptx / .ppt")

    tmp_path = None
    try:
        if suffix == ".ppt":
            tmp_path = _convert_ppt_to_pptx(path)
            pptx_path = tmp_path
        else:
            pptx_path = path

        prs = Presentation(str(pptx_path))
        slides_text = []

        for i, slide in enumerate(prs.slides, start=1):
            lines = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            # 备注页（speaker notes）
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip() if notes_tf else ""
                if notes_text:
                    lines.append(f"[Notes] {notes_text}")
            slides_text.append("\n".join(lines))

        return "\n\n".join(slides_text)

    finally:
        if tmp_path and tmp_path.exists():
            tmp_path.unlink(missing_ok=True)


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("用法: python extract_ppt.py <file.pptx>")
        sys.exit(1)
    print(extract_ppt(sys.argv[1]))